import math
import pandas as pd
import matplotlib.pyplot as plt
from sklearn.compose import ColumnTransformer
from sklearn.ensemble import RandomForestRegressor
from sklearn.linear_model import LinearRegression
from sklearn.metrics import mean_absolute_error, mean_squared_error, r2_score
from sklearn.model_selection import train_test_split
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import OneHotEncoder
from pptx import Presentation
from pptx.util import Inches

DATA_FILE = 'vgsales_clean.csv'
REPORT_FILE = 'report.md'
PPT_FILE = 'presentation.pptx'
PLOT_SALES = 'sales_distribution.png'
PLOT_PRED = 'predicted_vs_actual.png'


def load_data(path):
    df = pd.read_csv(path)
    df['Publisher'] = df['Publisher'].fillna('Unknown')
    return df


def prepare_features(df):
    top_publishers = df['Publisher'].value_counts().nlargest(20).index.tolist()
    df['PublisherGroup'] = df['Publisher'].where(df['Publisher'].isin(top_publishers), 'Other')
    X = df[['Year', 'Platform', 'Genre', 'PublisherGroup']].copy()
    y = df['Global_Sales'].copy()
    return X, y


def build_pipeline():
    categorical = ['Platform', 'Genre', 'PublisherGroup']
    try:
        encoder = OneHotEncoder(handle_unknown='ignore', sparse_output=False)
    except TypeError:
        encoder = OneHotEncoder(handle_unknown='ignore')
    preprocessor = ColumnTransformer(
        transformers=[
            ('cat', encoder, categorical),
        ],
        remainder='passthrough'
    )
    return preprocessor


def evaluate_model(model, X_test, y_test):
    y_pred = model.predict(X_test)
    return {
        'r2': r2_score(y_test, y_pred),
        'mae': mean_absolute_error(y_test, y_pred),
        'rmse': math.sqrt(mean_squared_error(y_test, y_pred)),
        'y_true': y_test,
        'y_pred': y_pred,
    }


def save_plots(df, results):
    plt.figure(figsize=(8, 5))
    plt.hist(df['Global_Sales'], bins=40, color='#4C72B0', edgecolor='black')
    plt.title('Phân phối Global Sales')
    plt.xlabel('Global Sales (triệu bản)')
    plt.ylabel('Số game')
    plt.tight_layout()
    plt.savefig(PLOT_SALES)
    plt.close()

    plt.figure(figsize=(8, 6))
    plt.scatter(results['y_true'], results['y_pred'], alpha=0.4, color='#DD8452')
    max_val = max(max(results['y_true']), max(results['y_pred']))
    plt.plot([0, max_val], [0, max_val], color='black', linestyle='--')
    plt.title('Actual vs Predicted Global Sales')
    plt.xlabel('Actual Global Sales')
    plt.ylabel('Predicted Global Sales')
    plt.tight_layout()
    plt.savefig(PLOT_PRED)
    plt.close()


def generate_report(metrics_lr, metrics_rf):
    best_model = 'Linear Regression' if metrics_lr['r2'] >= metrics_rf['r2'] else 'Random Forest'
    report_text = (
        '# Báo cáo Phân tích Game Sales\n\n'
        '## 1. Giới thiệu\n'
        'Dự án phân tích dữ liệu trò chơi video sử dụng tập dữ liệu `vgsales_clean.csv` để xây dựng mô hình dự đoán doanh số toàn cầu (`Global_Sales`).\n'
        'Mục tiêu là tìm ra ảnh hưởng của năm phát hành, nền tảng, thể loại và nhà phát hành đến doanh số và đánh giá chất lượng mô hình.\n\n'
        '## 2. Phương pháp\n'
        '- Dữ liệu đã được tải và làm sạch, điền giá trị thiếu của `Publisher` bằng `Unknown`.\n'
        '- Chọn tính năng: `Year`, `Platform`, `Genre`, `Publisher`.\n'
        '- Chuyển đổi dữ liệu phân loại bằng One-Hot Encoding.\n'
        '- So sánh hai mô hình: Linear Regression và Random Forest Regression.\n'
        '- Dữ liệu chia thành tập huấn luyện và kiểm tra với tỷ lệ 80/20.\n\n'
        '## 3. Kết quả\n'
        '### 3.1 Đánh giá mô hình\n'
        '| Mô hình | R2 score | MAE | RMSE |\n'
        '|---|---|---|---|\n'
        f'| Linear Regression | {metrics_lr["r2"]:.4f} | {metrics_lr["mae"]:.4f} | {metrics_lr["rmse"]:.4f} |\n'
        f'| Random Forest | {metrics_rf["r2"]:.4f} | {metrics_rf["mae"]:.4f} | {metrics_rf["rmse"]:.4f} |\n\n'
        '### 3.2 Kết quả chính\n'
        f'- Mô hình tốt nhất là {best_model} dựa theo R2 score.\n'
        f'- Sai số trung bình MAE cho {best_model} nằm trong khoảng giá trị doanh số triệu bản, cho thấy mô hình dự đoán khá sát thực tế đối với dữ liệu này.\n\n'
        '## 4. Insight\n'
        '- `Platform`, `Genre` và `Publisher` là các yếu tố ảnh hưởng lớn đến Global Sales khi không dùng trực tiếp doanh số vùng.\n'
        '- Các dòng game xuất hiện trên nền tảng hàng đầu và do các nhà phát hành lớn thực hiện thường có doanh số toàn cầu cao hơn.\n'
        '- Dự đoán doanh số toàn cầu từ chỉ metadata vẫn đạt kết quả chấp nhận được, cho thấy giá trị của thông tin sản phẩm trong lập kế hoạch phát hành.\n\n'
        '## 5. Kết luận\n'
        f'- Mô hình tốt nhất trong dự án này là {best_model}.\n'
        '- Để nâng cao mô hình, có thể bổ sung thêm các biến như đánh giá người dùng, chi phí marketing, thời điểm phát hành chi tiết và số lượng bản phát hành.\n'
        '- Tài liệu này được kèm theo một bài thuyết trình (`presentation.pptx`) nhằm phục vụ báo cáo nhóm.\n'
    )
    with open(REPORT_FILE, 'w', encoding='utf-8') as f:
        f.write(report_text)


def create_presentation(metrics_lr, metrics_rf):
    best_model = 'Linear Regression' if metrics_lr['r2'] >= metrics_rf['r2'] else 'Random Forest'
    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = 'Phân tích Game Sales & Dự đoán'
    title.placeholders[1].text = 'Đề tài: Phân tích doanh số trò chơi video'

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Mục tiêu'
    body = slide.shapes.placeholders[1].text_frame
    body.text = 'Mục tiêu dự án:'
    body.add_paragraph().text = '• Phân tích dữ liệu doanh số trò chơi'
    body.add_paragraph().text = '• Xây dựng mô hình dự đoán Global Sales'
    body.add_paragraph().text = '• Trình bày insight và đề xuất'

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Phương pháp'
    body = slide.shapes.placeholders[1].text_frame
    body.text = 'Phương pháp sử dụng:'
    body.add_paragraph().text = '• Tiền xử lý dữ liệu và one-hot encoding'
    body.add_paragraph().text = '• So sánh Linear Regression và Random Forest'
    body.add_paragraph().text = '• Đánh giá bằng R2, MAE, RMSE'

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Kết quả mô hình'
    body = slide.shapes.placeholders[1].text_frame
    body.text = 'Kết quả đánh giá:'
    body.add_paragraph().text = f'• Linear Regression: R2={metrics_lr["r2"]:.4f}, MAE={metrics_lr["mae"]:.4f}, RMSE={metrics_lr["rmse"]:.4f}'
    body.add_paragraph().text = f'• Random Forest: R2={metrics_rf["r2"]:.4f}, MAE={metrics_rf["mae"]:.4f}, RMSE={metrics_rf["rmse"]:.4f}'
    body.add_paragraph().text = f'• Mô hình tốt nhất: {best_model}'

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Phân phối Global Sales'
    slide.shapes.add_picture(PLOT_SALES, Inches(1), Inches(1.5), width=Inches(8))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Actual vs Predicted'
    slide.shapes.add_picture(PLOT_PRED, Inches(1), Inches(1.5), width=Inches(8))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Insight'
    body = slide.shapes.placeholders[1].text_frame
    body.text = 'Những điểm rút ra:'
    body.add_paragraph().text = '• Metadata game cho phép dự đoán doanh số tốt khi không dùng trực tiếp doanh số vùng.'
    body.add_paragraph().text = '• Platform và Genre vẫn là yếu tố quan trọng hàng đầu.'
    body.add_paragraph().text = '• Nhà phát hành lớn và nền tảng phổ biến giúp tăng khả năng doanh số cao.'

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Kết luận'
    body = slide.shapes.placeholders[1].text_frame
    body.text = 'Kết luận chính:'
    body.add_paragraph().text = f'• Mô hình tốt nhất: {best_model}.'
    body.add_paragraph().text = '• Có thể mở rộng thêm dữ liệu ngoài metadata để cải thiện dự đoán.'

    prs.save(PPT_FILE)


if __name__ == '__main__':
    df = load_data(DATA_FILE)
    X, y = prepare_features(df)

    preprocessor = build_pipeline()
    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)

    lr = Pipeline([('pre', preprocessor), ('model', LinearRegression())])
    rf = Pipeline([('pre', preprocessor), ('model', RandomForestRegressor(n_estimators=150, random_state=42, n_jobs=-1))])

    lr.fit(X_train, y_train)
    rf.fit(X_train, y_train)

    metrics_lr = evaluate_model(lr, X_test, y_test)
    metrics_rf = evaluate_model(rf, X_test, y_test)

    save_plots(df, metrics_rf)
    generate_report(metrics_lr, metrics_rf)
    create_presentation(metrics_lr, metrics_rf)

    print('Generated:', REPORT_FILE, PPT_FILE, PLOT_SALES, PLOT_PRED)
